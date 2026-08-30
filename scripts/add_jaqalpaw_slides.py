"""Add the "analytic protocol -> JaqalPaw" slides to docs/SpinBoson.pptx.

The three slides land directly after the protocol slide (page 6), which already
shows the analytic pulse as Δ(t), ϕ(t), g(t): convert it, show the JaqalPaw
pulse that results, then note the one part that does not fit on two tones.

Re-running replaces the slides it added last time rather than appending
duplicates: every slide it creates carries a marker shape name, and any slide
holding that marker is dropped before the new ones go in.

Regenerate the figures first:

    julia --project=. scripts/export_analytic_jaqalpaw.jl
    .venv/bin/python jaqal/verify_waveform.py --dump results/data/emulated_waveform.json
    julia --project=. scripts/plot_jaqalpaw_export.jl

Then:

    .venv/bin/python scripts/add_jaqalpaw_slides.py

Close the deck in PowerPoint first — it holds the file open and would
overwrite these edits on its next save.
"""

import os
import shutil
import sys
import tempfile

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

DECK = "docs/SpinBoson.pptx"
FIGDIR = "results/figures"
MARKER = "JaqalPawSlideMarker"

# The protocol slide these follow on from (1-based page number).
AFTER_PAGE = 6

BLANK_LAYOUT = 6

# Geometry, matched to the existing slides: a bold title near the top edge and
# the figure filling everything below it.
TITLE_LEFT, TITLE_TOP = Inches(0.61), Inches(0.20)
TITLE_W, TITLE_H = Inches(12.1), Inches(0.55)
TITLE_PT = 28
BODY_TOP = Inches(0.88)
BODY_H = Inches(6.45)

INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0xB0, 0x1C, 0x2E)
BLUE = RGBColor(0x1F, 0x4E, 0x9C)


def drop_previous(prs):
    """Remove slides a previous run added, identified by the marker shape.

    Dropping the relationship unlinks a slide but leaves its part in the
    package, so the caller must round-trip through `reload_clean` before adding
    replacements — otherwise the new slides claim the same part names as the
    orphans and the .pptx ends up with duplicate zip entries.
    """
    id_list = prs.slides._sldIdLst
    entries = list(id_list)
    removed = 0
    # Walk backwards so the surviving indices stay valid as entries are pulled.
    for idx in reversed(range(len(entries))):
        if any(sh.name == MARKER for sh in prs.slides[idx].shapes):
            prs.part.drop_rel(entries[idx].rId)
            id_list.remove(entries[idx])
            removed += 1
    return removed


def reload_clean(prs):
    """Serialise and re-open, so unlinked parts are gone for good.

    Saving walks the relationship graph, so the orphans left behind by
    `drop_previous` are not written; reloading then gives a package whose part
    names match what is actually in the deck.
    """
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        path = tmp.name
    try:
        prs.save(path)
        return Presentation(path)
    finally:
        os.unlink(path)


def move_after(prs, first_new, after_page):
    """Move the slides appended from index `first_new` onward so they land
    directly after `after_page` (1-based). python-pptx can only append, so the
    reordering happens on the slide-id list afterwards."""
    id_list = prs.slides._sldIdLst
    entries = list(id_list)
    new = entries[first_new:]
    for e in new:
        id_list.remove(e)
    for offset, e in enumerate(new):
        id_list.insert(after_page + offset, e)


def add_slide(prs, title):
    """Append a blank slide with the deck's standard bold title."""
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_W, TITLE_H)
    box.name = MARKER
    tf = box.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = title
    run.font.size = Pt(TITLE_PT)
    run.font.bold = True
    run.font.color.rgb = INK
    return slide


def add_figure(slide, slide_width, path, top=BODY_TOP, max_h=BODY_H,
               max_w=Inches(12.6)):
    """Place a PNG centred below the title, scaled to fit."""
    with Image.open(path) as im:
        px_w, px_h = im.size
    h = max_h
    w = Emu(int(h * px_w / px_h))
    if w > max_w:
        w = max_w
        h = Emu(int(w * px_h / px_w))
    left = Emu(int((slide_width - w) / 2))
    slide.shapes.add_picture(path, left, top, width=w, height=h)


def add_text(slide, left, top, width, height, blocks, size=15, gap=Pt(6)):
    """A text box built from (text, style) blocks.

    Styles: 'h' section heading, 'b' body line, 'eq' equation line,
    'note' small muted line.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, style in blocks:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.space_after = gap
        run = para.add_run()
        run.text = text
        f = run.font
        if style == "h":
            f.size, f.bold, f.color.rgb = Pt(size + 2), True, BLUE
            para.space_before = Pt(10)
        elif style == "eq":
            f.size, f.name, f.color.rgb = Pt(size), "Consolas", ACCENT
        elif style == "note":
            f.size, f.italic, f.color.rgb = Pt(size - 2), True, MUTED
        else:
            f.size, f.color.rgb = Pt(size), INK
    return box


def figure_path(name):
    path = os.path.join(FIGDIR, name)
    if not os.path.exists(path):
        raise SystemExit(
            f"Missing {path} — run scripts/plot_jaqalpaw_export.jl first."
        )
    return path


def build(prs):
    """The three slides that follow the protocol slide.

    The story is: slide 6 already shows the analytic protocol as Δ(t), ϕ(t),
    g(t) — so here we (1) convert it, (2) show the JaqalPaw pulse that comes
    out, (3) note the one thing that does not fit on two tones.
    """
    # --- How we convert ----------------------------------------------------
    s = add_slide(prs, "How we convert: Δ, ϕ, g → two sideband tones")
    add_text(
        s, TITLE_LEFT, Inches(0.80), Inches(12.4), Inches(0.95),
        [
            ("H = g(t)·a·[Jx·e^(−iΔt) + Jy·e^(+iΔt)·e^(−iϕ)] + h.c."
             "   =   [A(t)·a + B(t)·a†]·σ+ + h.c.,   σ+ = Jx + iJy", "eq"),
            ("A = (g/2)[e^(−iΔt) − i·e^(+i(Δt−ϕ))]  →  red tone,  rate 2|A|, "
             "phase −arg A        "
             "B = (g/2)[e^(+iΔt) − i·e^(−i(Δt−ϕ))]  →  blue tone,  rate 2|B|, "
             "phase −arg B", "eq"),
        ],
        size=13, gap=Pt(3),
    )
    add_figure(s, prs.slide_width, figure_path("jaqalpaw_tone_map.png"),
               top=Inches(1.85), max_h=Inches(5.5))

    # --- The JaqalPaw pulse ------------------------------------------------
    s = add_slide(prs, "The JaqalPaw pulse")
    add_text(
        s, TITLE_LEFT, Inches(0.80), Inches(12.4), Inches(1.5),
        [
            ("PulseData(GLOBAL_BEAM, 225.676e-6, freq0=200.000e6, amp0=100, "
             "phase0=0)", "eq"),
            ("PulseData(Q0, 225.676e-6,", "eq"),
            ("          freq0=232.100e6,  amp0=[512 values],  "
             "phase0=[512 values],     # blue sideband, tone 0", "eq"),
            ("          freq1=227.900e6,  amp1=[512 values],  "
             "phase1=[512 values])     # red sideband,  tone 1", "eq"),
            ("Static frequencies, 512 steps of 441 ns; compiles to 279 × "
             "256-bit words. Peak amplitude 24/100 at η = 0.1.", "note"),
        ],
        size=12, gap=Pt(1),
    )
    add_figure(s, prs.slide_width, figure_path("jaqalpaw_hardware.png"),
               top=Inches(2.55), max_h=Inches(4.8))

    # --- The caveat --------------------------------------------------------
    s = add_slide(prs, "Caveat: it is really a four-tone drive")
    add_figure(s, prs.slide_width, figure_path("jaqalpaw_four_tone.png"))


if __name__ == "__main__":
    lock = os.path.join(os.path.dirname(DECK), "~$" + os.path.basename(DECK))
    if os.path.exists(lock):
        sys.exit(
            f"{DECK} is open in PowerPoint ({lock} exists). Close it first — "
            "otherwise PowerPoint overwrites these edits on its next save."
        )

    backup = DECK.replace(".pptx", "_backup.pptx")
    if not os.path.exists(backup):
        shutil.copy2(DECK, backup)
        print(f"backed up original to {backup}")

    prs = Presentation(DECK)
    n_before = len(prs.slides._sldIdLst)
    removed = drop_previous(prs)
    if removed:
        prs = reload_clean(prs)
        print(f"replaced {removed} slide(s) from a previous run")

    n_existing = len(prs.slides._sldIdLst)
    build(prs)
    move_after(prs, n_existing, AFTER_PAGE)
    prs.save(DECK)
    print(f"{DECK}: {n_before} slides → {len(prs.slides._sldIdLst)}")
    for i, slide in enumerate(prs.slides):
        if any(sh.name == MARKER for sh in slide.shapes):
            title = next(sh.text_frame.text for sh in slide.shapes
                         if sh.name == MARKER)
            print(f"  slide {i + 1}: {title}")
